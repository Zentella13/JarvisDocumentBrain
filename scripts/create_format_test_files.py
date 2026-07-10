import sys
from pathlib import Path


def create_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 72), "El documento PDF menciona el adhesivo Novacote 5503.")
    doc.save(path)
    doc.close()


def create_word(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph("El documento Word contiene información sobre la máquina Miraflex.")
    document.save(path)


def create_excel(path: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet1 = workbook.active
    sheet1.title = "Materiales"
    sheet1["A1"] = "PET 12 micras"
    sheet1["A2"] = "BOPP HS"
    sheet1["A3"] = "CPP"

    sheet2 = workbook.create_sheet(title="Pruebas")
    sheet2["A1"] = "COF"
    sheet2["A2"] = "Seal Strength"
    sheet2["A3"] = "Thickness"

    workbook.save(path)


def create_powerpoint(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Fallas de sellado"
    body = slide.placeholders[1]
    body.text = "Presentación relacionada con fallas de sellado y laminación."

    slide2 = presentation.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Laminación"
    body2 = slide2.placeholders[1]
    body2.text = "Presentación relacionada con fallas de sellado y laminación."

    presentation.save(path)


def create_markdown(path: Path) -> None:
    path.write_text("# Procedimiento\nEste procedimiento menciona ASTM D882.\n", encoding="utf-8")


def main() -> int:
    output_dir = Path(r"C:\JarvisFormatTest")
    try:
        output_dir.mkdir(parents=True, exist_ok=True)
    except Exception as exc:
        print(f"Error al crear la carpeta de salida {output_dir}: {exc}")
        return 1

    creators = [
        (create_pdf, output_dir / "prueba_pdf.pdf"),
        (create_word, output_dir / "prueba_word.docx"),
        (create_excel, output_dir / "prueba_excel.xlsx"),
        (create_powerpoint, output_dir / "prueba_powerpoint.pptx"),
        (create_markdown, output_dir / "prueba_markdown.md"),
    ]

    errors = False
    for creator, path in creators:
        try:
            creator(path)
            print(f"Archivo creado: {path}")
        except Exception as exc:
            print(f"Error creando {path.name}: {exc}")
            errors = True

    if errors:
        print("Algunos archivos no se pudieron crear correctamente.")
        return 1

    print("Todos los archivos de prueba se crearon correctamente en", output_dir)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
