import os
from pathlib import Path
from typing import Any

import docx
import markdown
import pandas as pd
import pptx
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "md",
    ".markdown": "md",
}


def extract_text(file_path: str | os.PathLike[str]) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")

    if suffix == ".md" or suffix == ".markdown":
        content = path.read_text(encoding="utf-8", errors="ignore")
        return markdown.markdown(content)

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if suffix == ".docx":
        document = docx.Document(str(path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)

    if suffix == ".xlsx":
        excel = pd.ExcelFile(path)
        texts: list[str] = []
        for sheet_name in excel.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet_name)
            texts.append(df.astype(str).to_string(index=False))
        return "\n".join(texts)

    if suffix == ".pptx":
        presentation = pptx.Presentation(str(path))
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    raise ValueError(f"Unsupported extension: {suffix}")
